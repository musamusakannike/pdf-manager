import fitz  # PyMuPDF
from pypdf import PdfReader, PdfWriter
import os
import subprocess
from pdf2docx import Converter
import pandas as pd
import pdfplumber
from pptx import Presentation
from pptx.util import Inches
import io

class PDFEngine:
    """Core engine for PDF manipulations using PyMuPDF and pypdf."""
    
    @staticmethod
    def merge_pdfs(paths, output_path):
        writer = PdfWriter()
        for path in paths:
            reader = PdfReader(path)
            for page in reader.pages:
                writer.add_page(page)
        
        with open(output_path, "wb") as f:
            writer.write(f)
        return True

    @staticmethod
    def get_page_count(path):
        doc = fitz.open(path)
        count = doc.page_count
        doc.close()
        return count

    @staticmethod
    def render_page(path, page_num, zoom=2.0):
        doc = fitz.open(path)
        page = doc.load_page(page_num)
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat)
        
        # Convert to bytes for Qt
        img_data = pix.tobytes("ppm")
        doc.close()
        return img_data

    @staticmethod
    def split_pdf(path, output_dir):
        """Splits a PDF into individual pages."""
        reader = PdfReader(path)
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
            
        base_name = os.path.splitext(os.path.basename(path))[0]
        for i, page in enumerate(reader.pages):
            writer = PdfWriter()
            writer.add_page(page)
            output_path = os.path.join(output_dir, f"{base_name}_page_{i+1}.pdf")
            with open(output_path, "wb") as f:
                writer.write(f)
        return True

    @staticmethod
    def remove_pages(path, pages_to_keep, output_path):
        """Removes pages from a PDF by keeping only specified indices."""
        reader = PdfReader(path)
        writer = PdfWriter()
        for i in pages_to_keep:
            if i < len(reader.pages):
                writer.add_page(reader.pages[i])
        
        with open(output_path, "wb") as f:
            writer.write(f)
        return True

    @staticmethod
    def add_text_annotation(path, page_num, text, x, y, output_path):
        """Adds a text annotation to a specific page."""
        doc = fitz.open(path)
        page = doc.load_page(page_num)
        page.insert_text((x, y), text, color=(1, 0, 0)) # Red text
        doc.save(output_path)
        doc.close()
        return True

    @staticmethod
    def compress_pdf(path, output_path):
        """Simple compression by re-saving with optimization."""
        doc = fitz.open(path)
        doc.save(output_path, garbage=4, deflate=True)
        doc.close()
        return True

    @staticmethod
    def pdf_to_images(path, output_dir, fmt="png"):
        """Converts all PDF pages to image files."""
        doc = fitz.open(path)
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
            
        base_name = os.path.splitext(os.path.basename(path))[0]
        for i in range(len(doc)):
            page = doc.load_page(i)
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            pix.save(os.path.join(output_dir, f"{base_name}_{i+1}.{fmt}"))
        doc.close()
        return True
    
    @staticmethod
    def add_watermark(path, watermark_text, output_path, opacity=0.3):
        """Adds a watermark to all pages of a PDF."""
        doc = fitz.open(path)
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            rect = page.rect
            
            # Calculate center position
            x = rect.width / 2
            y = rect.height / 2
            
            # Use insert_text directly with a simple approach
            # Create text at center of page
            fontsize = 60
            font = fitz.Font("helv")
            text_length = font.text_length(watermark_text, fontsize=fontsize)
            
            # Position text at center
            text_point = fitz.Point(x - text_length/2, y)
            
            # Insert the text directly on the page with transparency
            tw = fitz.TextWriter(rect)
            tw.append(text_point, watermark_text, font=font, fontsize=fontsize)
            tw.write_text(page, color=(0.7, 0.7, 0.7), opacity=opacity)
        
        doc.save(output_path)
        doc.close()
        return True
    
    @staticmethod
    def encrypt_pdf(path, output_path, user_password, owner_password=None):
        """Encrypts a PDF with password protection."""
        doc = fitz.open(path)
        if owner_password is None:
            owner_password = user_password
        
        perm = int(
            fitz.PDF_PERM_ACCESSIBILITY |
            fitz.PDF_PERM_PRINT |
            fitz.PDF_PERM_COPY |
            fitz.PDF_PERM_ANNOTATE
        )
        
        encrypt_meth = fitz.PDF_ENCRYPT_AES_256
        doc.save(output_path, encryption=encrypt_meth, 
                user_pw=user_password, owner_pw=owner_password, permissions=perm)
        doc.close()
        return True
    
    @staticmethod
    def decrypt_pdf(path, output_path, password):
        """Removes password protection from a PDF."""
        try:
            doc = fitz.open(path)
            if doc.is_encrypted:
                if not doc.authenticate(password):
                    return False
            doc.save(output_path, encryption=fitz.PDF_ENCRYPT_NONE)
            doc.close()
            return True
        except:
            return False
    
    @staticmethod
    def extract_text(path, page_num=None):
        """Extracts text from a PDF page or entire document."""
        doc = fitz.open(path)
        if page_num is not None:
            page = doc.load_page(page_num)
            text = page.get_text()
        else:
            text = ""
            for page in doc:
                text += page.get_text()
        doc.close()
        return text
    
    @staticmethod
    def get_pdf_info(path):
        """Gets metadata information about a PDF."""
        doc = fitz.open(path)
        info = {
            'pages': len(doc),
            'title': doc.metadata.get('title', 'N/A'),
            'author': doc.metadata.get('author', 'N/A'),
            'subject': doc.metadata.get('subject', 'N/A'),
            'creator': doc.metadata.get('creator', 'N/A'),
            'producer': doc.metadata.get('producer', 'N/A'),
            'creation_date': doc.metadata.get('creationDate', 'N/A'),
            'modification_date': doc.metadata.get('modDate', 'N/A'),
            'encrypted': doc.is_encrypted
        }
        doc.close()
        return info
    
    @staticmethod
    def rotate_pages(path, output_path, rotation=90, pages=None):
        """Rotates specified pages or all pages by given degrees (90, 180, 270)."""
        doc = fitz.open(path)
        if pages is None:
            pages = range(len(doc))
        
        for page_num in pages:
            if page_num < len(doc):
                page = doc.load_page(page_num)
                page.set_rotation(rotation)
        
        doc.save(output_path)
        doc.close()
        return True

    @staticmethod
    def pdf_to_word(path, output_path):
        """Converts PDF to Word (.docx) using pdf2docx."""
        try:
            cv = Converter(path)
            cv.convert(output_path, start=0, end=None)
            cv.close()
            return True
        except Exception as e:
            print(f"Error converting PDF to Word: {e}")
            return False

    @staticmethod
    def word_to_pdf(path, output_path):
        """Converts Word (.docx) to PDF using LibreOffice."""
        try:
            # LibreOffice headless conversion
            output_dir = os.path.dirname(output_path)
            subprocess.run([
                'libreoffice', '--headless', '--convert-to', 'pdf',
                '--outdir', output_dir, path
            ], check=True)
            
            # LibreOffice saves with same basename, ensure output_path matches
            # If output_path name is different from input basename + .pdf, rename it
            base_name = os.path.splitext(os.path.basename(path))[0]
            created_pdf = os.path.join(output_dir, f"{base_name}.pdf")
            
            if created_pdf != output_path and os.path.exists(created_pdf):
                if os.path.exists(output_path):
                    os.remove(output_path)
                os.rename(created_pdf, output_path)
                
            return os.path.exists(output_path)
        except Exception as e:
            print(f"Error converting Word to PDF: {e}")
            return False

    @staticmethod
    def pdf_to_excel(path, output_path):
        """Converts PDF tables to Excel (.xlsx) using pdfplumber."""
        try:
            with pdfplumber.open(path) as pdf:
                all_tables = []
                for page in pdf.pages:
                    tables = page.extract_tables()
                    for table in tables:
                        df = pd.DataFrame(table)
                        # Clean up: remove empty rows/cols if needed
                        all_tables.append(df)
                
                if not all_tables:
                    # No tables found
                    return False
                
                with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
                    for i, df in enumerate(all_tables):
                        df.to_excel(writer, sheet_name=f"Table_{i+1}", index=False, header=False)
            return True
        except Exception as e:
            print(f"Error converting PDF to Excel: {e}")
            return False

    @staticmethod
    def excel_to_pdf(path, output_path):
        """Converts Excel (.xlsx) to PDF using LibreOffice."""
        try:
            output_dir = os.path.dirname(output_path)
            subprocess.run([
                'libreoffice', '--headless', '--convert-to', 'pdf',
                '--outdir', output_dir, path
            ], check=True)
            
            base_name = os.path.splitext(os.path.basename(path))[0]
            created_pdf = os.path.join(output_dir, f"{base_name}.pdf")
            
            if created_pdf != output_path and os.path.exists(created_pdf):
                if os.path.exists(output_path):
                    os.remove(output_path)
                os.rename(created_pdf, output_path)
                
            return os.path.exists(output_path)
        except Exception as e:
            print(f"Error converting Excel to PDF: {e}")
            return False

    @staticmethod
    def pdf_to_powerpoint(path, output_path):
        """Converts PDF pages to PowerPoint slides (as images)."""
        try:
            doc = fitz.open(path)
            prs = Presentation()
            
            # Use 16:9 aspect ratio usually, but let's match PDF page size roughly
            # Default slide width is 10 inches, height is 7.5 inches
            
            for i in range(len(doc)):
                page = doc.load_page(i)
                
                # Render page to image
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2)) # High res
                img_data = pix.tobytes("png")
                image_stream = io.BytesIO(img_data)
                
                # Create blank slide
                result = prs.slides.add_slide(prs.slide_layouts[6]) # 6 is blank layout
                
                # Add image to slide, fitting to slide size
                # We need to calculate aspect ratios to fit properly
                # For simplicity, we'll just fill the slide or center it.
                # Let's try to fit the image centered
                
                left = top = Inches(0)
                # width = Inches(10)
                # height = Inches(7.5)
                # result.shapes.add_picture(image_stream, left, top, width=width)
                
                # Better approach: Just add picture at 0,0 and let it size naturally or fit
                result.shapes.add_picture(image_stream, left, top, height=Inches(7.5))
                
            prs.save(output_path)
            doc.close()
            return True
        except Exception as e:
            print(f"Error converting PDF to PowerPoint: {e}")
            return False

    @staticmethod
    def powerpoint_to_pdf(path, output_path):
        """Converts PowerPoint (.pptx) to PDF using LibreOffice."""
        try:
            output_dir = os.path.dirname(output_path)
            subprocess.run([
                'libreoffice', '--headless', '--convert-to', 'pdf',
                '--outdir', output_dir, path
            ], check=True)
            
            base_name = os.path.splitext(os.path.basename(path))[0]
            created_pdf = os.path.join(output_dir, f"{base_name}.pdf")
            
            if created_pdf != output_path and os.path.exists(created_pdf):
                if os.path.exists(output_path):
                    os.remove(output_path)
                os.rename(created_pdf, output_path)
                
            return os.path.exists(output_path)
        except Exception as e:
            print(f"Error converting PowerPoint to PDF: {e}")
            return False
